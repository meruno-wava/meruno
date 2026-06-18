#!/usr/bin/env python3
"""
Dairy Meruno — Sistem Informasi Arsip Pribadi Rumah Sakit
Flask Backend: parsing Excel & PDF dengan analisis cerdas
GitHub: https://github.com/JIAkbar/meruno
"""

from flask import Flask, request, redirect, send_from_directory, jsonify, Response
import os, json, re, datetime, math
from pathlib import Path
from io import BytesIO

app = Flask(__name__)

def safe_num(v):
    """Konversi nilai numerik ke float JSON-safe (NaN/Inf → 0)."""
    try:
        f = float(v)
        return 0 if (math.isnan(f) or math.isinf(f)) else round(f, 2)
    except: return 0

BASE = Path(__file__).parent
DATA_DIR = BASE / "data"
UPLOAD_DIR = BASE / "uploads"
DATA_DIR.mkdir(exist_ok=True)
UPLOAD_DIR.mkdir(exist_ok=True)

EXCEL_DATA_FILE = DATA_DIR / "excel_data.json"
PDF_DATA_FILE   = DATA_DIR / "pdf_data.json"
ARCHIVE_FILE    = DATA_DIR / "archive.json"

# ── Casemix Konstanta ─────────────────────────────────────────────────────────
CASEMIX_DIR      = BASE / "uploads" / "casemix"
CASEMIX_DATA_DIR = DATA_DIR / "casemix"
CASEMIX_DIR.mkdir(parents=True, exist_ok=True)
CASEMIX_DATA_DIR.mkdir(parents=True, exist_ok=True)

CASEMIX_SECTIONS = [
    'klaim', 'regulasi', 'spo', 'pengorganisasian',
    'pelayanan', 'kebijakan', 'ur', 'galeri', 'hotnews',
]
DOC_SECTIONS = ['regulasi', 'spo', 'pengorganisasian', 'pelayanan', 'kebijakan']
CASEMIX_ACCEPT = {
    'klaim':            ['.xlsx', '.xls'],
    'regulasi':         ['.pdf', '.docx', '.doc'],
    'spo':              ['.pdf', '.docx', '.doc'],
    'pengorganisasian': ['.pdf', '.docx', '.doc'],
    'pelayanan':        ['.pdf', '.docx', '.doc'],
    'kebijakan':        ['.pdf', '.docx', '.doc'],
    'ur':               ['.pdf', '.pptx', '.ppt'],
    'galeri':           ['.jpg', '.jpeg', '.png', '.webp'],
    'hotnews':          [],
}

# ── Supabase Storage ──────────────────────────────────────────────────────────
SUPABASE_URL = os.environ.get("SUPABASE_URL", "")
SUPABASE_KEY = os.environ.get("SUPABASE_KEY", "")
_sb = None

def get_supabase():
    global _sb
    if _sb is None and SUPABASE_URL and SUPABASE_KEY:
        from supabase import create_client
        _sb = create_client(SUPABASE_URL, SUPABASE_KEY)
    return _sb

def sb_upload(bucket: str, path: str, data: bytes, content_type: str = "application/octet-stream"):
    """Upload file ke Supabase Storage. Return public URL atau None kalau gagal."""
    sb = get_supabase()
    if not sb:
        return None
    try:
        sb.storage.from_(bucket).upload(path, data, {"content-type": content_type, "upsert": "true"})
        return f"{SUPABASE_URL}/storage/v1/object/public/{bucket}/{path}"
    except Exception as e:
        print(f"[Supabase upload error] {e}")
        return None

def sb_delete(bucket: str, path: str) -> bool:
    """Hapus file dari Supabase Storage."""
    sb = get_supabase()
    if not sb:
        return False
    try:
        sb.storage.from_(bucket).remove([path])
        return True
    except Exception as e:
        print(f"[Supabase delete error] {e}")
        return False

def sb_read_json(bucket: str, path: str):
    """Baca JSON file dari Supabase Storage."""
    sb = get_supabase()
    if not sb:
        return None
    try:
        data = sb.storage.from_(bucket).download(path)
        return json.loads(data.decode())
    except Exception:
        return None

def sb_write_json(bucket: str, path: str, obj: dict) -> bool:
    """Simpan JSON file ke Supabase Storage."""
    data = json.dumps(obj, ensure_ascii=False, indent=2).encode()
    url = sb_upload(bucket, path, data, "application/json")
    return url is not None

# ── File Compression ──────────────────────────────────────────────────────────

def compress_pdf(data: bytes) -> bytes:
    """Kompres PDF dengan pikepdf. Return compressed bytes atau original kalau gagal."""
    try:
        import pikepdf
        with pikepdf.open(BytesIO(data)) as pdf:
            out = BytesIO()
            pdf.save(out, compress_streams=True, object_stream_mode=pikepdf.ObjectStreamMode.generate)
            compressed = out.getvalue()
            # Hanya pakai kalau lebih kecil
            return compressed if len(compressed) < len(data) else data
    except Exception as e:
        print(f"[compress_pdf] {e}")
        return data

def compress_image(data: bytes, max_width: int = 1920, quality: int = 82) -> bytes:
    """Kompres gambar dengan Pillow. Return compressed bytes atau original kalau gagal."""
    try:
        from PIL import Image
        img = Image.open(BytesIO(data))
        fmt = img.format or "JPEG"
        # Resize kalau lebih besar dari max_width
        if img.width > max_width:
            ratio = max_width / img.width
            img = img.resize((max_width, int(img.height * ratio)), Image.LANCZOS)
        # Convert RGBA → RGB kalau format JPEG
        if fmt == "JPEG" and img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        out = BytesIO()
        if fmt == "PNG":
            img.save(out, format="PNG", optimize=True)
        else:
            img.save(out, format="JPEG", quality=quality, optimize=True)
        compressed = out.getvalue()
        return compressed if len(compressed) < len(data) else data
    except Exception as e:
        print(f"[compress_image] {e}")
        return data

def compress_pptx(data: bytes, img_quality: int = 75) -> bytes:
    """Kompres gambar di dalam PPTX dengan python-pptx + Pillow. Return compressed bytes."""
    try:
        from pptx import Presentation
        from PIL import Image
        prs = Presentation(BytesIO(data))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    img_part = shape.image
                    img_data = img_part.blob
                    try:
                        img = Image.open(BytesIO(img_data))
                        # Resize kalau lebih dari 1920px
                        if img.width > 1920:
                            ratio = 1920 / img.width
                            img = img.resize((1920, int(img.height * ratio)), Image.LANCZOS)
                        out = BytesIO()
                        if img.mode in ("RGBA", "P"):
                            img = img.convert("RGB")
                        img.save(out, format="JPEG", quality=img_quality, optimize=True)
                        # Ganti image blob
                        img_part._blob = out.getvalue()
                    except Exception:
                        continue
        out = BytesIO()
        prs.save(out)
        compressed = out.getvalue()
        return compressed if len(compressed) < len(data) else data
    except Exception as e:
        print(f"[compress_pptx] {e}")
        return data

# ── Konfigurasi AI (opsional) ────────────────────────────────────────────────
# Set environment variable ANTHROPIC_API_KEY untuk mengaktifkan AI summary
CLAUDE_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")

def ai_summarize(text, doc_type, filename):
    """Gunakan Claude AI jika API key tersedia, fallback ke analisis lokal."""
    if not CLAUDE_API_KEY:
        return None
    try:
        import anthropic
        client = anthropic.Anthropic(api_key=CLAUDE_API_KEY)
        msg = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=600,
            messages=[{"role": "user", "content":
                f"Analisis dokumen arsip rumah sakit '{filename}' (jenis: {doc_type}).\n"
                f"Buat ringkasan dalam Bahasa Indonesia:\n"
                f"1. Ringkasan singkat isi dokumen (2-3 kalimat)\n"
                f"2. Informasi penting yang ditemukan\n"
                f"3. Rekomendasi pengarsipan\n\n"
                f"Konten:\n{text[:3000]}"
            }]
        )
        return msg.content[0].text
    except Exception:
        return None


# ── Analisis dokumen lokal ────────────────────────────────────────────────────

def detect_doc_type(filename, text=""):
    fn, tx = filename.lower(), text.lower()
    patterns = {
        "Rekam Medis":          ["rekam medis", "anamnesis", "diagnosa", "rm-", "pemeriksaan fisik"],
        "Surat Rujukan":        ["surat rujukan", "rujukan", "ditujukan kepada rs"],
        "Laporan Administrasi": ["laporan", "rekapitulasi", "rekap bulanan", "administrasi"],
        "Data Pasien":          ["pasien", "nama pasien", "tgl lahir", "no rm", "no. rm"],
        "Dokumen BPJS":         ["bpjs", "jkn", "klaim bpjs", "peserta aktif"],
        "Inventaris":           ["inventaris", "stok obat", "persediaan", "alkes"],
        "Keuangan RS":          ["tagihan", "biaya perawatan", "pembayaran", "invoice rs"],
        "Sertifikat/SK":        ["sertifikat", "surat keputusan", "sk no", "akreditasi"],
    }
    for dtype, kws in patterns.items():
        if any(k in fn or k in tx for k in kws):
            return dtype
    return "Dokumen Arsip RS"

def extract_keywords(text, n=10):
    stop = {"yang","dan","di","ke","dari","dengan","untuk","ini","itu","pada","adalah",
            "telah","akan","atau","dalam","tidak","juga","sudah","oleh","sebagai","nama",
            "nomor","tanggal","tgl","the","of","in","to","is","and","or"}
    words = re.findall(r'\b[a-zA-Z]{4,}\b', text)
    freq = {}
    for w in words:
        wl = w.lower()
        if wl not in stop:
            freq[wl] = freq.get(wl, 0) + 1
    return [w for w, _ in sorted(freq.items(), key=lambda x: x[1], reverse=True)[:n]]

def simple_analysis(text, filename):
    dates = re.findall(r'\d{1,2}[/-]\d{1,2}[/-]\d{2,4}|\d{4}-\d{2}-\d{2}', text)[:5]
    ids   = re.findall(r'(?:RM|No|No\.|Nomor|REF)[.\s:-]*([A-Z0-9/-]{3,20})', text, re.IGNORECASE)[:5]
    return {
        "doc_type":    detect_doc_type(filename, text),
        "keywords":    extract_keywords(text),
        "dates":       dates,
        "identifiers": ids,
        "word_count":  len(text.split()),
        "char_count":  len(text),
    }

def save_archive(ftype, filename, doc_type, info="", doc_id="", sb_url=""):
    # Coba baca arsip dari Supabase dulu, fallback ke local
    archive = sb_read_json("uploads", "meta/archive.json") or []
    if not archive and ARCHIVE_FILE.exists():
        with open(ARCHIVE_FILE, encoding="utf-8") as f:
            try: archive = json.load(f)
            except: pass
    entry = {
        "id": doc_id,
        "type": ftype, "filename": filename,
        "doc_type": doc_type, "info": info,
        "uploaded_at": datetime.datetime.now().isoformat(),
    }
    if sb_url:
        entry["sb_url"] = sb_url
    archive.insert(0, entry)
    archive = archive[:200]
    # Simpan ke local
    with open(ARCHIVE_FILE, "w", encoding="utf-8") as f:
        json.dump(archive, f, ensure_ascii=False, indent=2)
    # Simpan ke Supabase
    sb_write_json("uploads", "meta/archive.json", archive)


# ── Routes ────────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return send_from_directory(str(BASE), "index.html")

@app.route("/excel.html")
def excel_page():
    return send_from_directory(str(BASE), "excel.html")

@app.route("/pdf.html")
def pdf_page():
    return send_from_directory(str(BASE), "pdf.html")

@app.route("/changelog.html")
def changelog_page():
    return send_from_directory(str(BASE), "changelog.html")

@app.route("/checklist.html")
def checklist_page():
    """Generate checklist.html langsung dari checklist.md."""
    md = BASE / "checklist.md"
    if not md.exists():
        return "checklist.md tidak ditemukan", 404
    with open(md, encoding="utf-8") as f:
        return Response(render_checklist(f.read()), mimetype="text/html")

@app.route("/<path:filename>")
def static_files(filename):
    fp = BASE / filename
    if fp.exists() and fp.is_file():
        return send_from_directory(str(BASE), filename)
    return f"404 — {filename} tidak ditemukan", 404


# ── Upload API ────────────────────────────────────────────────────────────────

@app.route("/api/upload", methods=["POST"])
def upload():
    if "file" not in request.files:
        return jsonify({"error": "Tidak ada file dikirim"}), 400
    f = request.files["file"]
    if not f.filename:
        return jsonify({"error": "Nama file kosong"}), 400
    fn = f.filename.lower()
    if fn.endswith((".xlsx", ".xls")):
        return process_excel(f)
    elif fn.endswith(".pdf"):
        return process_pdf(f)
    return jsonify({"error": f"Format tidak didukung: {fn}"}), 400


def process_excel(file):
    try:
        import openpyxl, pandas as pd
    except ImportError:
        return jsonify({"error": "Jalankan: pip install openpyxl pandas"}), 500
    try:
        fb = file.read()
        wb = openpyxl.load_workbook(BytesIO(fb), data_only=True)
        sheets = []

        for sn in wb.sheetnames:
            ws = wb[sn]
            hdr_row = next(ws.iter_rows(min_row=1, max_row=1, values_only=True))
            hdrs = [str(v) if v is not None else f"Kolom {i+1}" for i, v in enumerate(hdr_row)]
            rows = []
            for row in ws.iter_rows(min_row=2, values_only=True):
                r = []
                for v in row:
                    if isinstance(v, (datetime.datetime, datetime.date)):
                        r.append(str(v)[:10])
                    elif v is None:
                        r.append("")
                    else:
                        r.append(str(v))
                rows.append(r)
                if len(rows) >= 500:
                    break

            df = pd.read_excel(BytesIO(fb), sheet_name=sn, nrows=500)
            num_cols = df.select_dtypes(include="number").columns.tolist()
            col_stats = {}
            for c in num_cols[:8]:
                try:
                    col_stats[c] = {
                        "sum":   round(float(df[c].sum()), 2),
                        "avg":   round(float(df[c].mean()), 2),
                        "min":   round(float(df[c].min()), 2),
                        "max":   round(float(df[c].max()), 2),
                        "count": int(df[c].count()),
                    }
                except: pass

            sheets.append({
                "name": sn, "headers": hdrs, "rows": rows,
                "row_count": len(rows), "col_count": len(hdrs),
                "numeric_cols": num_cols, "col_stats": col_stats,
            })

        all_text = " ".join(" ".join(r) for s in sheets for r in s["rows"][:50])
        ana = simple_analysis(all_text, file.filename)
        ai  = ai_summarize(all_text[:3000], ana["doc_type"], file.filename)

        doc_id = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        data = {
            "id":          doc_id,
            "filename":    file.filename,
            "uploaded_at": datetime.datetime.now().strftime("%d %B %Y, %H:%M"),
            "doc_type":    ana["doc_type"],
            "keywords":    ana["keywords"],
            "sheets":      sheets,
            "ai_summary":  ai,
        }
        # Upload file ke Supabase Storage
        sb_url = sb_upload("uploads", f"archive/{doc_id}_{file.filename}", fb,
                           "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        id_file = DATA_DIR / f"excel_{doc_id}.json"
        with open(id_file, "w", encoding="utf-8") as f2:
            json.dump(data, f2, ensure_ascii=False, indent=2)
        with open(EXCEL_DATA_FILE, "w", encoding="utf-8") as f2:
            json.dump(data, f2, ensure_ascii=False, indent=2)
        # Simpan JSON data ke Supabase
        sb_write_json("uploads", f"meta/excel_{doc_id}.json", data)
        save_archive("excel", file.filename, ana["doc_type"],
                     f"{sum(s['row_count'] for s in sheets)} baris", doc_id, sb_url or "")
        return redirect("/excel.html")
    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "detail": traceback.format_exc()}), 500


def process_pdf(file):
    try:
        import pdfplumber
    except ImportError:
        return jsonify({"error": "Jalankan: pip install pdfplumber"}), 500
    try:
        fb = file.read()
        fb = compress_pdf(fb)
        pages, full_text = [], ""

        with pdfplumber.open(BytesIO(fb)) as pdf:
            for i, pg in enumerate(pdf.pages):
                txt  = pg.extract_text() or ""
                tbls = []
                for t in (pg.extract_tables() or []):
                    if t:
                        tbls.append([[str(c) if c else "" for c in row] for row in t])
                full_text += txt + "\n"
                pages.append({
                    "page_num": i + 1, "text": txt,
                    "tables": tbls, "char_count": len(txt),
                })

        ana = simple_analysis(full_text, file.filename)
        ai  = ai_summarize(full_text, ana["doc_type"], file.filename)

        if not ai:
            # Fallback ringkasan sederhana
            first_lines = [l.strip() for l in full_text.split("\n") if l.strip()][:6]
            ai = "Cuplikan konten dokumen:\n" + "\n".join(f"• {l}" for l in first_lines)

        doc_id = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        data = {
            "id":          doc_id,
            "filename":    file.filename,
            "uploaded_at": datetime.datetime.now().strftime("%d %B %Y, %H:%M"),
            "page_count":  len(pages),
            "doc_type":    ana["doc_type"],
            "keywords":    ana["keywords"],
            "dates":       ana.get("dates", []),
            "identifiers": ana.get("identifiers", []),
            "ai_summary":  ai,
            "pages":       pages,
        }
        # Upload file ke Supabase Storage
        sb_url = sb_upload("uploads", f"archive/{doc_id}_{file.filename}", fb, "application/pdf")
        id_file = DATA_DIR / f"pdf_{doc_id}.json"
        with open(id_file, "w", encoding="utf-8") as f2:
            json.dump(data, f2, ensure_ascii=False, indent=2)
        with open(PDF_DATA_FILE, "w", encoding="utf-8") as f2:
            json.dump(data, f2, ensure_ascii=False, indent=2)
        # Simpan JSON data ke Supabase
        sb_write_json("uploads", f"meta/pdf_{doc_id}.json", data)
        save_archive("pdf", file.filename, ana["doc_type"], f"{len(pages)} halaman", doc_id, sb_url or "")
        return redirect("/pdf.html")
    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "detail": traceback.format_exc()}), 500


# ── Data APIs ─────────────────────────────────────────────────────────────────

@app.route("/api/data/excel")
def get_excel():
    doc_id = request.args.get("id", "")
    fp = (DATA_DIR / f"excel_{doc_id}.json") if doc_id else EXCEL_DATA_FILE
    if fp.exists():
        with open(fp, encoding="utf-8") as f:
            return jsonify(json.load(f))
    return jsonify(None)

@app.route("/api/data/pdf")
def get_pdf():
    doc_id = request.args.get("id", "")
    fp = (DATA_DIR / f"pdf_{doc_id}.json") if doc_id else PDF_DATA_FILE
    if fp.exists():
        with open(fp, encoding="utf-8") as f:
            return jsonify(json.load(f))
    return jsonify(None)

@app.route("/api/archive")
def get_archive():
    ftype = request.args.get("type", "")
    # Coba Supabase dulu
    data = sb_read_json("uploads", "meta/archive.json")
    # Fallback ke local file
    if data is None and ARCHIVE_FILE.exists():
        with open(ARCHIVE_FILE, encoding="utf-8") as f:
            try: data = json.load(f)
            except: data = None
    if data is None:
        data = []
    if ftype:
        data = [d for d in data if d.get("type") == ftype]
    return jsonify(data)

@app.route("/api/delete/<ftype>/<doc_id>", methods=["DELETE"])
def delete_doc(ftype, doc_id):
    if ftype not in ("excel", "pdf"):
        return jsonify({"error": "Tipe tidak valid"}), 400
    # Hapus file data spesifik
    id_file = DATA_DIR / f"{ftype}_{doc_id}.json"
    deleted = False
    if id_file.exists():
        id_file.unlink()
        deleted = True
    # Hapus dari arsip
    if ARCHIVE_FILE.exists():
        with open(ARCHIVE_FILE, encoding="utf-8") as f:
            try: archive = json.load(f)
            except: archive = []
        archive = [a for a in archive if not (a.get("type") == ftype and a.get("id") == doc_id)]
        with open(ARCHIVE_FILE, "w", encoding="utf-8") as f:
            json.dump(archive, f, ensure_ascii=False, indent=2)
    # Update file "current" jika yang dihapus adalah yang aktif
    main_file = EXCEL_DATA_FILE if ftype == "excel" else PDF_DATA_FILE
    if main_file.exists():
        with open(main_file, encoding="utf-8") as f:
            try: cur = json.load(f)
            except: cur = {}
        if cur.get("id") == doc_id:
            remaining = sorted(DATA_DIR.glob(f"{ftype}_*.json"), key=lambda p: p.name, reverse=True)
            if remaining:
                with open(remaining[0], encoding="utf-8") as f:
                    next_data = json.load(f)
                with open(main_file, "w", encoding="utf-8") as f:
                    json.dump(next_data, f, ensure_ascii=False, indent=2)
            else:
                main_file.unlink()
    return jsonify({"ok": True, "deleted": deleted})


# ── Checklist renderer ────────────────────────────────────────────────────────

def render_checklist(md_content):
    sections, current = [], None
    for line in md_content.split("\n"):
        line = line.rstrip()
        if line.startswith("## "):
            if current: sections.append(current)
            current = {"title": line[3:].strip(), "items": []}
        elif line.startswith("- [") and current:
            done = line.startswith("- [x]")
            txt  = re.sub(r'^- \[.\] ', '', line).strip()
            txt  = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', txt)
            txt  = re.sub(r'`(.+?)`', r'<code>\1</code>', txt)
            current["items"].append({"done": done, "text": txt})
    if current: sections.append(current)

    total      = sum(len(s["items"]) for s in sections)
    done_count = sum(1 for s in sections for i in s["items"] if i["done"])
    pct        = round(done_count / total * 100) if total else 0

    items_html = ""
    for sec in sections:
        items_html += f'<div class="cl-sec"><div class="cl-sec-title">{sec["title"]}</div>'
        for item in sec["items"]:
            cls = "done" if item["done"] else "pending"
            box = "✓" if item["done"] else ""
            items_html += (f'<div class="chk-item {cls}">'
                           f'<div class="chk-box {cls}">{box}</div>'
                           f'<div class="chk-label">{item["text"]}</div></div>')
        items_html += "</div>"

    return f"""<!DOCTYPE html>
<html lang="id">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>Checklist — Dairy Meruno</title>
<link href="https://fonts.googleapis.com/css2?family=Playfair+Display:wght@700&family=Nunito+Sans:wght@400;600;700&family=Space+Mono&display=swap" rel="stylesheet">
<style>
:root{{--teal:#0d6e7a;--teal-dark:#094e58;--orange:#e8821e;--green:#1f9068;--green-light:#27b880;--cream:#f5f0e8;--cream-dark:#ede5d8;--dark:#122830;--gray:#6b8a92;--white:#fff}}
*{{box-sizing:border-box;margin:0;padding:0}}body{{font-family:'Nunito Sans',sans-serif;background:var(--cream);color:var(--dark)}}
.topbar{{background:var(--teal-dark);padding:14px 5%;display:flex;align-items:center;justify-content:space-between}}
.logo{{font-family:'Playfair Display',serif;font-size:18px;color:white;text-decoration:none;font-weight:700}}
.nav{{display:flex;gap:24px}}.nav a{{color:rgba(255,255,255,.7);text-decoration:none;font-size:13px;font-weight:600;transition:.2s}}.nav a:hover{{color:white}}
.hero{{background:linear-gradient(135deg,var(--teal-dark),#1a3d47);padding:50px 5%;text-align:center}}
.hero-badge{{display:inline-block;background:rgba(255,255,255,.1);border:1px solid rgba(255,255,255,.15);border-radius:50px;padding:5px 16px;font-size:11px;font-weight:700;letter-spacing:2px;text-transform:uppercase;color:var(--orange);margin-bottom:16px}}
.hero h1{{font-family:'Playfair Display',serif;font-size:36px;color:white;margin-bottom:8px}}
.hero p{{color:rgba(255,255,255,.6);font-size:15px}}
.progress-wrap{{max-width:560px;margin:22px auto 0}}
.progress-label{{display:flex;justify-content:space-between;font-size:13px;color:rgba(255,255,255,.7);margin-bottom:8px}}
.progress-bar{{height:10px;background:rgba(255,255,255,.15);border-radius:5px;overflow:hidden}}
.progress-fill{{height:100%;border-radius:5px;background:linear-gradient(90deg,var(--green-light),#1a8fa0);width:{pct}%}}
.content{{max-width:820px;margin:0 auto;padding:40px 5% 80px}}
.back-btn{{display:inline-flex;align-items:center;gap:8px;background:var(--teal);color:white;text-decoration:none;padding:10px 22px;border-radius:50px;font-weight:700;font-size:14px;margin-bottom:28px;transition:.2s}}.back-btn:hover{{background:var(--teal-dark)}}
.stats-row{{display:flex;gap:14px;margin-bottom:28px;flex-wrap:wrap}}
.stat-card{{background:var(--white);border-radius:12px;border:1.5px solid var(--cream-dark);padding:18px 22px;flex:1;min-width:100px;text-align:center}}
.stat-card .num{{font-family:'Space Mono',monospace;font-size:30px;font-weight:700;color:var(--teal)}}
.stat-card.g .num{{color:var(--green)}}.stat-card.o .num{{color:var(--orange)}}
.stat-card .lbl{{font-size:11px;color:var(--gray);text-transform:uppercase;letter-spacing:1px;margin-top:4px}}
.cl-sec{{background:var(--white);border-radius:14px;border:1.5px solid var(--cream-dark);margin-bottom:18px;overflow:hidden}}
.cl-sec-title{{padding:13px 22px;background:var(--cream);border-bottom:1px solid var(--cream-dark);font-size:14px;font-weight:700}}
.chk-item{{display:flex;align-items:flex-start;gap:12px;padding:11px 20px;border-bottom:1px solid var(--cream-dark);transition:.15s}}.chk-item:last-child{{border-bottom:none}}.chk-item:hover{{background:var(--cream)}}
.chk-box{{width:20px;height:20px;border-radius:5px;display:flex;align-items:center;justify-content:center;font-size:11px;flex-shrink:0;margin-top:1px}}
.chk-box.done{{background:var(--green);color:white}}.chk-box.pending{{border:2px solid var(--gray);background:white}}
.chk-label{{font-size:14px;line-height:1.5}}.chk-item.done .chk-label{{text-decoration:line-through;color:var(--gray)}}
code{{font-family:'Space Mono',monospace;font-size:12px;background:var(--cream);padding:1px 6px;border-radius:4px;color:var(--teal)}}strong{{color:var(--teal-dark)}}
.note{{font-size:12px;color:var(--gray);text-align:center;margin-top:20px;padding:12px;background:var(--cream);border-radius:8px;border:1px solid var(--cream-dark)}}
</style>
</head>
<body>
<div class="topbar">
  <a href="/" class="logo">🗂️ Dairy Meruno</a>
  <nav class="nav">
    <a href="/">Beranda</a><a href="/excel.html">Excel</a>
    <a href="/pdf.html">PDF</a><a href="/changelog.html">Changelog</a>
    <a href="/checklist.html">Checklist</a>
  </nav>
</div>
<div class="hero">
  <div class="hero-badge">Progress Pengembangan</div>
  <h1>Checklist Fitur</h1>
  <p>Sistem Informasi Arsip Pribadi — Rumah Sakit</p>
  <div class="progress-wrap">
    <div class="progress-label"><span>Progress Total</span><span>{done_count}/{total} selesai ({pct}%)</span></div>
    <div class="progress-bar"><div class="progress-fill"></div></div>
  </div>
</div>
<div class="content">
  <a href="/" class="back-btn">← Kembali ke Beranda</a>
  <div class="stats-row">
    <div class="stat-card g"><div class="num">{done_count}</div><div class="lbl">Selesai</div></div>
    <div class="stat-card o"><div class="num">{total - done_count}</div><div class="lbl">Tersisa</div></div>
    <div class="stat-card"><div class="num">{total}</div><div class="lbl">Total</div></div>
    <div class="stat-card"><div class="num">{pct}%</div><div class="lbl">Progress</div></div>
  </div>
  {items_html}
  <p class="note">📝 Halaman ini di-generate otomatis dari <code>checklist.md</code> — edit file .md untuk memperbarui.</p>
</div>
</body></html>"""


# ── Casemix Routes ───────────────────────────────────────────────────────────

SECTION_TITLES = {
    'klaim':            'Rekapitulasi Klaim',
    'regulasi':         'Pedoman & Regulasi',
    'spo':              'SPO',
    'pengorganisasian': 'Pedoman Pengorganisasian',
    'pelayanan':        'Pedoman Pelayanan',
    'kebijakan':        'Kebijakan Layanan',
    'ur':               'Utilization Review',
    'galeri':           'Galeri & Memory',
    'hotnews':          'Hot News',
}

@app.route("/casemix")
def casemix_dashboard():
    return send_from_directory(str(BASE), "casemix.html")

@app.route("/casemix/<section>")
def casemix_section(section):
    if section not in CASEMIX_SECTIONS:
        return "Section tidak ditemukan", 404
    if section in DOC_SECTIONS:
        return send_from_directory(str(BASE), "casemix_doc.html")
    page_map = {
        'klaim':   'casemix_klaim.html',
        'ur':      'casemix_ur.html',
        'galeri':  'casemix_galeri.html',
        'hotnews': 'casemix_hotnews.html',
    }
    page = page_map.get(section)
    if page:
        fp = BASE / page
        if fp.exists():
            return send_from_directory(str(BASE), page)
        return f"Halaman {page} belum dibuat", 404
    return "Halaman belum tersedia", 404

@app.route("/casemix/api/files/<section>")
def casemix_files(section):
    if section not in CASEMIX_SECTIONS:
        return jsonify({"files": [], "error": "Section tidak valid"}), 400
    CASEMIX_DATA_DIR.mkdir(parents=True, exist_ok=True)
    # Coba Supabase dulu
    sb_data = sb_read_json("casemix", f"meta/{section}.json")
    if sb_data is not None:
        return jsonify(sb_data)
    # Fallback ke local file
    meta_file = CASEMIX_DATA_DIR / f"{section}.json"
    if meta_file.exists():
        with open(meta_file, encoding="utf-8") as fp:
            try:
                return jsonify(json.load(fp))
            except Exception:
                pass
    return jsonify({"files": [], "count": 0})

@app.route("/casemix/api/news")
def casemix_news():
    CASEMIX_DATA_DIR.mkdir(parents=True, exist_ok=True)
    # Coba Supabase dulu
    sb_data = sb_read_json("casemix", "meta/news.json")
    if sb_data is not None:
        return jsonify(sb_data)
    # Fallback ke local file
    news_file = CASEMIX_DATA_DIR / "news.json"
    if news_file.exists():
        with open(news_file, encoding="utf-8") as fp:
            try:
                return jsonify(json.load(fp))
            except Exception:
                pass
    return jsonify({"news": [], "count": 0})

@app.route("/casemix/api/sync-news", methods=["POST"])
def casemix_sync_news():
    """Sync berita terbaru dari wavahusada.com via WordPress REST API."""
    import urllib.request, html, re as _re
    WP_API = "https://wavahusada.com/wp-json/wp/v2/posts?categories=110&per_page=5&_embed&_fields=id,title,date,excerpt,link,_embedded"
    try:
        req = urllib.request.Request(WP_API, headers={"User-Agent": "DairyMeruno/1.0"})
        with urllib.request.urlopen(req, timeout=10) as resp:
            posts = json.loads(resp.read().decode())
    except Exception as e:
        return jsonify({"error": f"Gagal mengambil data: {e}"}), 502

    news_file = CASEMIX_DATA_DIR / "news.json"
    # Coba baca dari Supabase dulu, fallback ke local
    existing = sb_read_json("casemix", "meta/news.json")
    if existing is None:
        existing = json.loads(news_file.read_text(encoding="utf-8")) if news_file.exists() else {"news": [], "count": 0}
    existing_ids = {n["id"] for n in existing["news"]}

    def strip_html(s):
        return _re.sub(r"<[^>]+>", "", html.unescape(s or "")).strip()

    def fmt_date(iso):
        try:
            dt = datetime.datetime.fromisoformat(iso)
            MONTHS = ["","Januari","Februari","Maret","April","Mei","Juni",
                      "Juli","Agustus","September","Oktober","November","Desember"]
            return f"{dt.day} {MONTHS[dt.month]} {dt.year}"
        except:
            return iso

    added = 0
    for p in posts:
        wp_id = f"wp_{p['id']}"
        if wp_id in existing_ids:
            continue
        img_url = None
        try:
            img_url = p["_embedded"]["wp:featuredmedia"][0]["source_url"]
        except (KeyError, IndexError, TypeError):
            pass
        entry = {
            "id": wp_id,
            "title": html.unescape(p["title"]["rendered"]),
            "content": strip_html(p["excerpt"]["rendered"]),
            "source_url": p["link"],
            "published_at": p["date"],
            "published_at_display": fmt_date(p["date"]),
            "img_url": img_url,
        }
        existing["news"].insert(0, entry)
        existing_ids.add(wp_id)
        added += 1

    existing["count"] = len(existing["news"])
    # Simpan ke local
    news_file.write_text(json.dumps(existing, ensure_ascii=False, indent=2), encoding="utf-8")
    # Simpan ke Supabase
    sb_write_json("casemix", "meta/news.json", existing)
    return jsonify({"added": added, "total": existing["count"]})

@app.route("/casemix/api/section-info")
def casemix_section_info():
    section = request.args.get("section", "")
    if not section or section not in CASEMIX_SECTIONS:
        return jsonify({"error": "Parameter section tidak valid"}), 400
    return jsonify({
        "title": SECTION_TITLES.get(section, section),
        "is_doc": section in DOC_SECTIONS,
        "accept": CASEMIX_ACCEPT.get(section, []),
    })

@app.route("/casemix/files/<section>/<filename>")
def casemix_file_serve(section, filename):
    # Kalau Supabase tersedia, redirect ke public URL
    if get_supabase() and SUPABASE_URL:
        sb_path = f"{section}/{filename}"
        sb_public_url = f"{SUPABASE_URL}/storage/v1/object/public/casemix/{sb_path}"
        return redirect(sb_public_url)
    # Fallback ke local file
    section_dir = CASEMIX_DIR / section
    fp = section_dir / filename
    if not fp.exists():
        return "File tidak ditemukan", 404
    return send_from_directory(str(section_dir), filename)

@app.route("/casemix/upload/<section>", methods=["POST"])
def casemix_upload(section):
    if section not in CASEMIX_SECTIONS:
        return jsonify({"error": "Section tidak valid"}), 400
    if section == 'hotnews':
        return casemix_add_news()
    if "file" not in request.files:
        return jsonify({"error": "Tidak ada file dikirim"}), 400
    f = request.files["file"]
    if not f.filename:
        return jsonify({"error": "Nama file kosong"}), 400
    fn = f.filename.lower()
    allowed = CASEMIX_ACCEPT.get(section, [])
    if allowed and not any(fn.endswith(ext) for ext in allowed):
        exts = ", ".join(allowed)
        return jsonify({"error": f"Format tidak didukung. Gunakan: {exts}"}), 400
    if section == 'klaim':
        return casemix_process_klaim(f)
    if section == 'galeri':
        return casemix_save_galeri(f)
    return casemix_save_doc(f, section)

def casemix_save_doc(file, section):
    section_dir = CASEMIX_DIR / section
    section_dir.mkdir(parents=True, exist_ok=True)
    doc_id = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_fn = re.sub(r'[^\w\-.]', '_', file.filename)
    stored_name = f"{doc_id}_{safe_fn}"
    fb = file.read()
    original_size = len(fb)
    ext = Path(file.filename).suffix.lower()
    if ext == ".pdf":
        fb = compress_pdf(fb)
    elif ext in (".pptx", ".ppt"):
        fb = compress_pptx(fb)
    title = request.form.get("title", file.filename)
    desc  = request.form.get("description", "")
    # Upload ke Supabase
    import mimetypes
    ct = mimetypes.guess_type(file.filename)[0] or "application/octet-stream"
    sb_url = sb_upload("casemix", f"{section}/{stored_name}", fb, ct)
    # Fallback ke local storage kalau Supabase tidak tersedia
    if not sb_url:
        with open(str(section_dir / stored_name), "wb") as lf:
            lf.write(fb)
    entry = {
        "id":               doc_id,
        "title":            title,
        "filename":         stored_name,
        "original_name":    file.filename,
        "description":      desc,
        "size":             len(fb),
        "uploaded_at":      datetime.datetime.now().isoformat(),
        "uploaded_at_display": datetime.datetime.now().strftime("%d %B %Y, %H:%M"),
        "file_url":         sb_url if sb_url else f"/casemix/files/{section}/{stored_name}",
        "section":          section,
    }
    # Baca metadata: coba Supabase dulu, fallback ke local
    meta = sb_read_json("casemix", f"meta/{section}.json")
    if meta is None:
        meta_file = CASEMIX_DATA_DIR / f"{section}.json"
        meta = {"files": [], "count": 0}
        if meta_file.exists():
            with open(meta_file, encoding="utf-8") as fp:
                try: meta = json.load(fp)
                except: pass
    meta["files"].insert(0, entry)
    meta["count"] = len(meta["files"])
    # Simpan metadata ke local dan Supabase
    meta_file = CASEMIX_DATA_DIR / f"{section}.json"
    with open(meta_file, "w", encoding="utf-8") as fp:
        json.dump(meta, fp, ensure_ascii=False, indent=2)
    sb_write_json("casemix", f"meta/{section}.json", meta)
    return jsonify({"ok": True, "id": doc_id, "entry": entry,
                    "original_size": original_size, "compressed_size": len(fb)})

def casemix_save_galeri(file):
    section_dir = CASEMIX_DIR / "galeri"
    section_dir.mkdir(parents=True, exist_ok=True)
    doc_id = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_fn = re.sub(r'[^\w\-.]', '_', file.filename)
    stored_name = f"{doc_id}_{safe_fn}"
    fb = file.read()
    original_size = len(fb)
    fb = compress_image(fb)
    caption = request.form.get("caption", "")
    # Upload ke Supabase
    import mimetypes
    ct = mimetypes.guess_type(file.filename)[0] or "image/jpeg"
    sb_url = sb_upload("casemix", f"galeri/{stored_name}", fb, ct)
    # Fallback ke local storage kalau Supabase tidak tersedia
    if not sb_url:
        with open(str(section_dir / stored_name), "wb") as lf:
            lf.write(fb)
    entry = {
        "id":               doc_id,
        "caption":          caption,
        "filename":         stored_name,
        "original_name":    file.filename,
        "uploaded_at":      datetime.datetime.now().isoformat(),
        "uploaded_at_display": datetime.datetime.now().strftime("%d %B %Y, %H:%M"),
        "img_url":          sb_url if sb_url else f"/casemix/files/galeri/{stored_name}",
        "size":             len(fb),
    }
    # Baca metadata: coba Supabase dulu, fallback ke local
    meta = sb_read_json("casemix", "meta/galeri.json")
    if meta is None:
        meta_file = CASEMIX_DATA_DIR / "galeri.json"
        meta = {"files": [], "count": 0}
        if meta_file.exists():
            with open(meta_file, encoding="utf-8") as fp:
                try: meta = json.load(fp)
                except: pass
    meta["files"].insert(0, entry)
    meta["count"] = len(meta["files"])
    # Simpan ke local dan Supabase
    meta_file = CASEMIX_DATA_DIR / "galeri.json"
    with open(meta_file, "w", encoding="utf-8") as fp:
        json.dump(meta, fp, ensure_ascii=False, indent=2)
    sb_write_json("casemix", "meta/galeri.json", meta)
    return jsonify({"ok": True, "id": doc_id, "entry": entry,
                    "original_size": original_size, "compressed_size": len(fb)})

def casemix_process_klaim(file):
    if file.filename.lower().endswith('.xls'):
        return jsonify({"error": "Format .xls lama tidak didukung. Buka di Excel → Save As → .xlsx dulu."}), 400
    try:
        import openpyxl, pandas as pd
    except ImportError:
        return jsonify({"error": "Jalankan: pip install openpyxl pandas"}), 500
    try:
        fb = file.read()
        original_size = len(fb)
        wb = openpyxl.load_workbook(BytesIO(fb), data_only=True)
        sheets = []
        for sn in wb.sheetnames:
            ws = wb[sn]
            hdr_row = next(ws.iter_rows(min_row=1, max_row=1, values_only=True))
            hdrs = [str(v) if v is not None else f"Kolom {i+1}" for i, v in enumerate(hdr_row)]
            rows = []
            for row in ws.iter_rows(min_row=2, values_only=True):
                r = []
                for v in row:
                    if isinstance(v, (datetime.datetime, datetime.date)):
                        r.append(str(v)[:10])
                    elif v is None:
                        r.append("")
                    else:
                        r.append(str(v))
                rows.append(r)
                if len(rows) >= 1000:
                    break
            df = pd.read_excel(BytesIO(fb), sheet_name=sn, nrows=1000)
            num_cols = df.select_dtypes(include="number").columns.tolist()
            col_stats = {}
            for c in num_cols[:12]:
                try:
                    col_stats[c] = {
                        "sum":   safe_num(df[c].sum()),
                        "avg":   safe_num(df[c].mean()),
                        "min":   safe_num(df[c].min()),
                        "max":   safe_num(df[c].max()),
                        "count": int(df[c].count()),
                    }
                except: pass
            sheets.append({
                "name":        sn,
                "headers":     hdrs,
                "rows":        rows,
                "row_count":   len(rows),
                "col_count":   len(hdrs),
                "numeric_cols": num_cols,
                "col_stats":   col_stats,
            })
        doc_id = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        entry = {
            "id":               doc_id,
            "title":            request.form.get("title", file.filename),
            "original_name":    file.filename,
            "uploaded_at":      datetime.datetime.now().isoformat(),
            "uploaded_at_display": datetime.datetime.now().strftime("%d %B %Y, %H:%M"),
            "sheets":           sheets,
            "total_rows":       sum(s["row_count"] for s in sheets),
            "size":             original_size,
        }
        # Scrape + delete: file TIDAK disimpan, hanya data JSON-nya
        # Baca metadata: coba Supabase dulu, fallback ke local
        meta = sb_read_json("casemix", "meta/klaim.json")
        if meta is None:
            meta_file = CASEMIX_DATA_DIR / "klaim.json"
            meta = {"files": [], "count": 0}
            if meta_file.exists():
                with open(meta_file, encoding="utf-8") as fp:
                    try: meta = json.load(fp)
                    except: pass
        meta["files"].insert(0, entry)
        meta["count"] = len(meta["files"])
        # Simpan ke local dan Supabase
        meta_file = CASEMIX_DATA_DIR / "klaim.json"
        with open(meta_file, "w", encoding="utf-8") as fp:
            json.dump(meta, fp, ensure_ascii=False, indent=2)
        sb_write_json("casemix", "meta/klaim.json", meta)
        return jsonify({"ok": True, "id": doc_id, "entry": entry, "count": meta["count"], "original_size": original_size, "compressed_size": original_size})
    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "detail": traceback.format_exc()}), 500

def casemix_add_news():
    title   = request.form.get("title", "").strip()
    content = request.form.get("content", "").strip()
    if not title:
        return jsonify({"error": "Judul wajib diisi"}), 400
    news_id = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    img_original_size = 0
    img_compressed_size = 0
    entry = {
        "id":                 news_id,
        "title":              title,
        "content":            content,
        "published_at":       datetime.datetime.now().isoformat(),
        "published_at_display": datetime.datetime.now().strftime("%d %B %Y, %H:%M"),
        "img_url":            None,
    }
    if "image" in request.files:
        img = request.files["image"]
        if img.filename:
            news_dir = CASEMIX_DIR / "hotnews"
            news_dir.mkdir(parents=True, exist_ok=True)
            safe_fn = re.sub(r'[^\w\-.]', '_', img.filename)
            stored = f"{news_id}_{safe_fn}"
            img_bytes = img.read()
            img_original_size = len(img_bytes)
            img_bytes = compress_image(img_bytes)
            img_compressed_size = len(img_bytes)
            # Upload gambar ke Supabase
            import mimetypes
            ct = mimetypes.guess_type(img.filename)[0] or "image/jpeg"
            sb_img_url = sb_upload("casemix", f"hotnews/{stored}", img_bytes, ct)
            if sb_img_url:
                entry["img_url"] = sb_img_url
            else:
                # Fallback: simpan lokal
                with open(str(news_dir / stored), "wb") as lf:
                    lf.write(img_bytes)
                entry["img_url"] = f"/casemix/files/hotnews/{stored}"
    # Baca metadata: coba Supabase dulu, fallback ke local
    data = sb_read_json("casemix", "meta/news.json")
    if data is None:
        news_file = CASEMIX_DATA_DIR / "news.json"
        data = {"news": [], "count": 0}
        if news_file.exists():
            with open(news_file, encoding="utf-8") as fp:
                try: data = json.load(fp)
                except: pass
    data["news"].insert(0, entry)
    data["count"] = len(data["news"])
    # Simpan ke local dan Supabase
    news_file = CASEMIX_DATA_DIR / "news.json"
    with open(news_file, "w", encoding="utf-8") as fp:
        json.dump(data, fp, ensure_ascii=False, indent=2)
    sb_write_json("casemix", "meta/news.json", data)
    return jsonify({"ok": True, "id": news_id, "entry": entry,
                    "original_size": img_original_size if img_original_size else 0,
                    "compressed_size": img_compressed_size if img_compressed_size else 0})

@app.route("/casemix/delete/<section>/<doc_id>", methods=["DELETE"])
def casemix_delete(section, doc_id):
    if section not in CASEMIX_SECTIONS:
        return jsonify({"error": "Section tidak valid"}), 400
    meta_json_key = "news.json" if section == "hotnews" else f"{section}.json"
    meta_file = CASEMIX_DATA_DIR / meta_json_key
    # Baca metadata: coba Supabase dulu, fallback ke local
    meta = sb_read_json("casemix", f"meta/{meta_json_key}")
    if meta is None:
        meta = {}
        if meta_file.exists():
            with open(meta_file, encoding="utf-8") as fp:
                try: meta = json.load(fp)
                except: pass
    list_key = "news" if section == "hotnews" else "files"
    items = meta.get(list_key, [])
    new_items = []
    deleted_entry = None
    for item in items:
        if item.get("id") == doc_id:
            deleted_entry = item
        else:
            new_items.append(item)
    meta[list_key] = new_items
    meta["count"] = len(new_items)
    # Simpan metadata yang diupdate ke local dan Supabase
    with open(meta_file, "w", encoding="utf-8") as fp:
        json.dump(meta, fp, ensure_ascii=False, indent=2)
    sb_write_json("casemix", f"meta/{meta_json_key}", meta)
    # Hapus file fisik jika ada (klaim tidak punya file fisik)
    if deleted_entry:
        fn = deleted_entry.get("filename") or deleted_entry.get("img_url", "").split("/")[-1]
        if fn:
            # Hapus dari Supabase Storage
            sb_delete("casemix", f"{section}/{fn}")
            # Hapus dari local storage jika ada
            fp2 = CASEMIX_DIR / section / fn
            if fp2.exists():
                fp2.unlink()
    return jsonify({"ok": True, "deleted": deleted_entry is not None})


# ── Main ──────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import webbrowser, threading

    print()
    print("=" * 55)
    print("  🗂️  DAIRY MERUNO — Sistem Arsip RS")
    print("  📍  http://localhost:8080")
    print("  ⚡  Flask " + __import__("flask").__version__)
    if CLAUDE_API_KEY:
        print("  🤖  AI aktif (Claude API)")
    else:
        print("  💡  Set ANTHROPIC_API_KEY untuk mengaktifkan AI")
    print("=" * 55)
    print()

    def open_browser():
        import time; time.sleep(1.5)
        webbrowser.open("http://localhost:8080")
    threading.Thread(target=open_browser, daemon=True).start()

    port = int(os.environ.get("PORT", 8080))
    host = "0.0.0.0" if os.environ.get("PORT") else "127.0.0.1"
    app.run(host=host, port=port, debug=False)
